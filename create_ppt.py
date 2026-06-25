import sys
sys.path.insert(0, 'D:\\pythonlibs')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY = RGBColor(0x1A, 0x3A, 0x6C)
SECONDARY = RGBColor(0x2E, 0x86, 0xDE)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2D, 0x34, 0x36)
GRAY_TEXT = RGBColor(0x6C, 0x75, 0x7D)
LIGHT_GRAY = RGBColor(0xE9, 0xEC, 0xEF)

def add_blank_slide():
    slide_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(slide_layout)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide_content(slide, items, left, top, width, height, font_size=16, color=DARK_TEXT, spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = Pt(font_size * 0.5)
        p.line_spacing = spacing
    return txBox

def add_decorative_header(slide, title, subtitle=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=PRIMARY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), Inches(13.333), Inches(0.06), fill_color=SECONDARY)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7), title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4), subtitle, font_size=14, color=RGBColor(0xB0, 0xC4, 0xDE))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(12.2), Inches(0.3), Inches(0.6), Inches(0.6), fill_color=ACCENT)

def add_footer(slide, page_num, total_pages):
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3), fill_color=LIGHT_GRAY)
    add_textbox(slide, Inches(0.8), Inches(7.22), Inches(8), Inches(0.25), '苏州晶讯科技股份有限公司保密办 | 商业秘密保护知识培训', font_size=10, color=GRAY_TEXT)
    add_textbox(slide, Inches(11.5), Inches(7.22), Inches(1.5), Inches(0.25), f'{page_num} / {total_pages}', font_size=10, color=GRAY_TEXT, alignment=PP_ALIGN.RIGHT)

TOTAL_PAGES = 18
current_page = 0

def next_page():
    global current_page
    current_page += 1
    return current_page

# ========== Slide 1: Cover ==========
slide = add_blank_slide()
set_bg(slide, PRIMARY)

# Decorative shapes
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5), fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.333), Inches(0.08), fill_color=SECONDARY)
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.5), Inches(0.8), Inches(2.5), Inches(2.5), fill_color=SECONDARY)
shape = add_shape(slide, MSO_SHAPE.OVAL, Inches(11), Inches(1.3), Inches(1.8), Inches(1.8), fill_color=ACCENT)

add_textbox(slide, Inches(1.2), Inches(2.2), Inches(10), Inches(1.2), '商业秘密保护知识培训', font_size=48, bold=True, color=WHITE)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.5), Inches(1.5), Inches(0.08), fill_color=ACCENT)
add_textbox(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(0.6), 'Trade Secret Protection Training', font_size=20, color=RGBColor(0xB0, 0xC4, 0xDE))
add_textbox(slide, Inches(1.2), Inches(5.2), Inches(10), Inches(0.5), '苏州晶讯科技股份有限公司保密办', font_size=20, color=WHITE)
add_textbox(slide, Inches(1.2), Inches(5.8), Inches(10), Inches(0.4), '2024年7月', font_size=16, color=RGBColor(0xB0, 0xC4, 0xDE))

# ========== Slide 2: Contents ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '目  录', 'CONTENTS')

chapters = [
    ('01', '基础认知', '认识商业秘密', '定义、特点、范围与内容'),
    ('02', '法律合规', '法律法规与责任', '相关法律、禁止行为、法律责任'),
    ('03', '风险识别', '窃密手段与泄密风险', '窃密方式、失密泄密情形'),
    ('04', '责任体系', '保密工作组织架构', '委员会、保密办、部门、员工职责'),
    ('05', '公司制度', '保密管理与处罚', '密级划分、标识规定、处罚措施'),
    ('06', '防护措施', '企业自我保护', '门禁、文件、人员、信息安全管理'),
]

for i, (num, title_cn, title_sub, desc) in enumerate(chapters):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(2.0 + row * 2.4)
    w = Inches(3.7)
    h = Inches(2.0)
    
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill_color=WHITE)
    card.adjustments[0] = 0.05
    
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), h, fill_color=SECONDARY)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.2), Inches(1.5), Inches(0.5), num, font_size=32, bold=True, color=ACCENT)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.75), Inches(3), Inches(0.4), title_cn, font_size=20, bold=True, color=PRIMARY)
    add_textbox(slide, left + Inches(0.3), top + Inches(1.15), Inches(3), Inches(0.35), title_sub, font_size=12, color=GRAY_TEXT)
    add_textbox(slide, left + Inches(0.3), top + Inches(1.5), Inches(3.2), Inches(0.4), desc, font_size=11, color=GRAY_TEXT)

add_footer(slide, next_page(), TOTAL_PAGES)

# ========== Slide 3: Chapter 1 ==========
slide = add_blank_slide()
set_bg(slide, PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5), fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.5), Inches(5), Inches(3), Inches(3), fill_color=SECONDARY)
add_shape(slide, MSO_SHAPE.OVAL, Inches(9.5), Inches(4), Inches(2), Inches(2), fill_color=ACCENT)

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(3), Inches(1), '01', font_size=100, bold=True, color=ACCENT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2), Inches(1), Inches(0.08), fill_color=ACCENT)
add_textbox(slide, Inches(1.5), Inches(4.5), Inches(8), Inches(0.8), '基础认知', font_size=40, bold=True, color=WHITE)
add_textbox(slide, Inches(1.5), Inches(5.3), Inches(8), Inches(0.5), '认识商业秘密 · 定义、特点、范围与内容', font_size=18, color=RGBColor(0xB0, 0xC4, 0xDE))

# ========== Slide 4: Definition & Features ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '商业秘密的定义与特点', 'Definition & Characteristics')
page = next_page()

# Definition card
def_card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.5), fill_color=WHITE)
def_card.adjustments[0] = 0.05
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(0.1), Inches(1.5), fill_color=ACCENT)
add_textbox(slide, Inches(1.2), Inches(1.85), Inches(2), Inches(0.4), '定义', font_size=18, bold=True, color=PRIMARY)
add_textbox(slide, Inches(1.2), Inches(2.3), Inches(11), Inches(0.9),
    '商业秘密——不为公众所知悉，能为权利人带来经济利益，具有实用性并经权利人采取保密措施的技术信息和经营信息。',
    font_size=16, color=DARK_TEXT)

# Four features
features = [
    ('秘密性', '不为公众所知悉'),
    ('价值性', '能带来经济利益'),
    ('实用性', '具有实际使用价值'),
    ('保密性', '权利人采取保密措施'),
]

for i, (title, desc) in enumerate(features):
    left = Inches(0.8 + i * 3.0)
    top = Inches(3.8)
    w = Inches(2.7)
    h = Inches(2.5)
    
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill_color=WHITE)
    card.adjustments[0] = 0.06
    
    icon_shape = add_shape(slide, MSO_SHAPE.OVAL, left + Inches(0.95), top + Inches(0.3), Inches(0.8), Inches(0.8), fill_color=SECONDARY)
    add_textbox(slide, left + Inches(0.95), top + Inches(0.42), Inches(0.8), Inches(0.6), str(i+1), font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, left, top + Inches(1.25), w, Inches(0.4), title, font_size=18, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), top + Inches(1.75), w - Inches(0.4), Inches(0.6), desc, font_size=12, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 5: Scope & Content ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '商业秘密的范围与内容', 'Scope & Content')
page = next_page()

# Two columns
col1_left = Inches(0.8)
col2_left = Inches(6.8)
col_w = Inches(5.7)
col_top = Inches(1.7)
col_h = Inches(5.0)

# Technical info
card1 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, col1_left, col_top, col_w, col_h, fill_color=WHITE)
card1.adjustments[0] = 0.03
add_shape(slide, MSO_SHAPE.RECTANGLE, col1_left, col_top, col_w, Inches(0.7), fill_color=SECONDARY)
add_textbox(slide, col1_left + Inches(0.3), col_top + Inches(0.15), col_w, Inches(0.45), '技术信息', font_size=20, bold=True, color=WHITE)

tech_items = [
    '设计方案',
    '程序代码',
    '产品配方',
    '制作工艺',
    '制作方法',
    '技术诀窍',
]
add_bullet_slide_content(slide, ['●  ' + item for item in tech_items],
    col1_left + Inches(0.5), col_top + Inches(1.0), col_w - Inches(1), Inches(3.5),
    font_size=15, color=DARK_TEXT)

# Business info
card2 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, col2_left, col_top, col_w, col_h, fill_color=WHITE)
card2.adjustments[0] = 0.03
add_shape(slide, MSO_SHAPE.RECTANGLE, col2_left, col_top, col_w, Inches(0.7), fill_color=ACCENT)
add_textbox(slide, col2_left + Inches(0.3), col_top + Inches(0.15), col_w, Inches(0.45), '经营信息', font_size=20, bold=True, color=WHITE)

biz_items = [
    '战略规划与管理方法',
    '商业模式与改制上市',
    '并购重组与产权交易',
    '财务信息与投融资决策',
    '产购销策略与资源储备',
    '客户信息与招投标事项',
]
add_bullet_slide_content(slide, ['●  ' + item for item in biz_items],
    col2_left + Inches(0.5), col_top + Inches(1.0), col_w - Inches(1), Inches(3.5),
    font_size=15, color=DARK_TEXT)

add_textbox(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.3),
    '依据：《中央企业商业秘密保护暂行规定》',
    font_size=12, color=GRAY_TEXT, alignment=PP_ALIGN.RIGHT)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 6: Chapter 2 ==========
slide = add_blank_slide()
set_bg(slide, PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5), fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.5), Inches(5), Inches(3), Inches(3), fill_color=SECONDARY)
add_shape(slide, MSO_SHAPE.OVAL, Inches(9.5), Inches(4), Inches(2), Inches(2), fill_color=ACCENT)

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(3), Inches(1), '02', font_size=100, bold=True, color=ACCENT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2), Inches(1), Inches(0.08), fill_color=ACCENT)
add_textbox(slide, Inches(1.5), Inches(4.5), Inches(8), Inches(0.8), '法律合规', font_size=40, bold=True, color=WHITE)
add_textbox(slide, Inches(1.5), Inches(5.3), Inches(8), Inches(0.5), '法律法规与责任 · 依法保护商业秘密', font_size=18, color=RGBColor(0xB0, 0xC4, 0xDE))

# ========== Slide 7: Laws & Regulations ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '相关法律法规', 'Laws & Regulations')
page = next_page()

laws = [
    ('反不正当竞争法', '第10条', '经营者不得以盗窃、利诱、胁迫等不正当手段侵犯他人商业秘密；不得披露、使用或允许他人使用以不正当手段获取的商业秘密；不得违反保密约定披露、使用或允许他人使用其所掌握的商业秘密。'),
    ('民法通则', '第118条', '公民、法人的著作权、专利权、商标专用权、发现权、发明权和其他科技成果权受到剽窃、篡改、假冒等侵害的，有权要求停止侵害、消除影响、赔偿损失。'),
    ('刑法', '第219条', '以不正当手段获取或违反保密约定披露、使用商业秘密，给权利人造成重大损失的，处三年以下有期徒刑，并处或单处罚金；造成特别严重后果的，处三年以上七年以下有期徒刑，并处罚金。'),
    ('劳动合同法', '第23/90条', '用人单位与劳动者可以在劳动合同中约定保守商业秘密和与知识产权相关的保密事项。劳动者违反保密义务或竞业限制，给用人单位造成损失的，应当承担赔偿责任。'),
]

for i, (law, article, content) in enumerate(laws):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.1)
    top = Inches(1.7 + row * 2.55)
    w = Inches(5.7)
    h = Inches(2.25)
    
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill_color=WHITE)
    card.adjustments[0] = 0.04
    
    tag = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.25), top + Inches(0.2), Inches(1.5), Inches(0.38), fill_color=ACCENT)
    tag.adjustments[0] = 0.3
    add_textbox(slide, left + Inches(0.25), top + Inches(0.24), Inches(1.5), Inches(0.3), article, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, left + Inches(2.0), top + Inches(0.22), Inches(3.5), Inches(0.4), law, font_size=17, bold=True, color=PRIMARY)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.75), w - Inches(0.6), Inches(1.4), content, font_size=12, color=DARK_TEXT)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 8: Prohibited Acts ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '法律禁止的侵犯商业秘密行为', 'Prohibited Unfair Competition Acts')
page = next_page()

acts = [
    ('01', '不正当获取', '以盗窃、利诱、胁迫或者其他不正当手段获取权利人的商业秘密'),
    ('02', '不当披露使用', '披露、使用或者允许他人使用以前项手段获取的权利人的商业秘密'),
    ('03', '违反约定披露', '与权利人有业务关系的单位和个人违反合同约定或保密要求，披露、使用或允许他人使用其所掌握的商业秘密'),
    ('04', '员工违规披露', '权利人的职工违反合同约定或保密要求，披露、使用或者允许他人使用其所掌握的商业秘密'),
    ('05', '第三人恶意获取', '第三人明知或者应知前几种侵犯商业秘密是违法行为，仍从那里获取、使用或者披露商业秘密'),
    ('06', '高薪挖人窃密', '以高薪或者其他优厚条件聘用掌握或者了解权利人商业秘密的人员，以获取、使用、披露商业秘密'),
]

for i, (num, title, desc) in enumerate(acts):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(1.7 + row * 2.55)
    w = Inches(3.7)
    h = Inches(2.25)
    
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill_color=WHITE)
    card.adjustments[0] = 0.04
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), h, fill_color=SECONDARY)
    
    add_textbox(slide, left + Inches(0.3), top + Inches(0.2), Inches(1), Inches(0.5), num, font_size=28, bold=True, color=ACCENT)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.8), w - Inches(0.5), Inches(0.4), title, font_size=16, bold=True, color=PRIMARY)
    add_textbox(slide, left + Inches(0.3), top + Inches(1.25), w - Inches(0.5), Inches(0.9), desc, font_size=11, color=GRAY_TEXT)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 9: Legal Liability ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '侵犯商业秘密的法律责任', 'Legal Liability')
page = next_page()

liabilities = [
    ('民事责任', '停止侵害、消除影响、赔礼道歉', 'SECONDARY'),
    ('赔偿责任', '给权利人造成损害的，承担损害赔偿责任', 'ACCENT'),
    ('费用承担', '承担权利人调查侵权所支付的合理费用', 'SECONDARY'),
    ('行政处罚', '根据情节处以一万元以上二十万元以下罚款', 'ACCENT'),
    ('刑事责任', '造成重大损失或特别严重后果的，依法追究刑事责任', 'PRIMARY'),
]

for i, (title, desc, color_key) in enumerate(liabilities):
    top = Inches(1.8 + i * 1.0)
    left = Inches(1.5)
    w = Inches(10.3)
    h = Inches(0.8)
    
    color = SECONDARY if color_key == 'SECONDARY' else (ACCENT if color_key == 'ACCENT' else PRIMARY)
    
    bar = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill_color=WHITE)
    bar.adjustments[0] = 0.1
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), h, fill_color=color)
    
    add_textbox(slide, left + Inches(0.5), top + Inches(0.22), Inches(2.5), Inches(0.4), title, font_size=17, bold=True, color=color)
    add_textbox(slide, left + Inches(3.2), top + Inches(0.25), Inches(7), Inches(0.4), desc, font_size=14, color=DARK_TEXT)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 10: Chapter 3 ==========
slide = add_blank_slide()
set_bg(slide, PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5), fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.5), Inches(5), Inches(3), Inches(3), fill_color=SECONDARY)
add_shape(slide, MSO_SHAPE.OVAL, Inches(9.5), Inches(4), Inches(2), Inches(2), fill_color=ACCENT)

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(3), Inches(1), '03', font_size=100, bold=True, color=ACCENT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2), Inches(1), Inches(0.08), fill_color=ACCENT)
add_textbox(slide, Inches(1.5), Inches(4.5), Inches(8), Inches(0.8), '风险识别', font_size=40, bold=True, color=WHITE)
add_textbox(slide, Inches(1.5), Inches(5.3), Inches(8), Inches(0.5), '窃密手段与泄密风险 · 提高警惕，防范未然', font_size=18, color=RGBColor(0xB0, 0xC4, 0xDE))

# ========== Slide 11: Espionage Methods ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '常用的窃密手段和方法', 'Common Espionage Methods')
page = next_page()

methods = [
    ('金钱收买', '用金钱诱惑相关人员'),
    ('色情勾引', '利用美色套取情报'),
    ('物质引诱', '以贵重物品为诱饵'),
    ('现场观察', '亲临现场直接观察'),
    ('攀拉关系', '交结朋友接近目标'),
    ('窃听窃照', '使用技术设备偷拍偷录'),
    ('抓柄要挟', '逼人把柄逼其就范'),
    ('报刊分析', '分析公开资料积累信息'),
    ('废品收集', '从废弃物品中筛选情报'),
    ('心战策反', '培养情感进行心理策反'),
]

for i, (title, desc) in enumerate(methods):
    row = i // 5
    col = i % 5
    left = Inches(0.6 + col * 2.5)
    top = Inches(1.8 + row * 2.55)
    w = Inches(2.2)
    h = Inches(2.2)
    
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill_color=WHITE)
    card.adjustments[0] = 0.08
    
    icon = add_shape(slide, MSO_SHAPE.OVAL, left + Inches(0.7), top + Inches(0.25), Inches(0.8), Inches(0.8), fill_color=SECONDARY)
    add_textbox(slide, left + Inches(0.7), top + Inches(0.38), Inches(0.8), Inches(0.6), str(i+1).zfill(2), font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, left, top + Inches(1.15), w, Inches(0.35), title, font_size=14, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), top + Inches(1.5), w - Inches(0.3), Inches(0.6), desc, font_size=10, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 12: Leak Risks ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '企业商业秘密泄露的主要途径', 'Main Leakage Channels')
page = next_page()

# Left: 失密
left_col = Inches(0.8)
right_col = Inches(6.8)
col_w = Inches(5.7)
col_top = Inches(1.7)
col_h = Inches(5.0)

card1 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_col, col_top, col_w, col_h, fill_color=WHITE)
card1.adjustments[0] = 0.03
add_shape(slide, MSO_SHAPE.RECTANGLE, left_col, col_top, col_w, Inches(0.7), fill_color=SECONDARY)
add_textbox(slide, left_col + Inches(0.3), col_top + Inches(0.15), col_w, Inches(0.45), '易发生失密的情况', font_size=18, bold=True, color=WHITE)
add_textbox(slide, left_col + Inches(0.3), col_top + Inches(0.75), col_w - Inches(0.6), Inches(0.4),
    '失密：即遗失或失掉了秘密。如秘密文件、资料、图纸以及个人使用的笔记本、优盘、移动硬盘等丢失。',
    font_size=11, color=GRAY_TEXT)

loss_items = [
    '保密制度不健全',
    '文件发放范围过大',
    '个人携带文件外出',
    '把秘密文件带回家看',
    '横向传阅文件',
    '安全措施不够',
]
add_bullet_slide_content(slide, ['●  ' + item for item in loss_items],
    left_col + Inches(0.5), col_top + Inches(1.5), col_w - Inches(1), Inches(3.2),
    font_size=13, color=DARK_TEXT)

# Right: 泄密
card2 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, right_col, col_top, col_w, col_h, fill_color=WHITE)
card2.adjustments[0] = 0.03
add_shape(slide, MSO_SHAPE.RECTANGLE, right_col, col_top, col_w, Inches(0.7), fill_color=ACCENT)
add_textbox(slide, right_col + Inches(0.3), col_top + Inches(0.15), col_w, Inches(0.45), '易发生泄密的情况', font_size=18, bold=True, color=WHITE)
add_textbox(slide, right_col + Inches(0.3), col_top + Inches(0.75), col_w - Inches(0.6), Inches(0.4),
    '泄密：将不宜外泄的内部秘密被人窃听或泄露给无关人员。包括口头、书面、技术文件等多种渠道。',
    font_size=11, color=GRAY_TEXT)

leak_items = [
    '公共场所谈论内情',
    '同乡好友聚会聊天',
    '为名求利，报道猎奇',
    '利害关系，丧失理智',
    '酒后失言吐露内情',
    '有求他人，讨好对方',
    '家人面前谈论商业秘密',
    '众人面前逞强好胜',
]
add_bullet_slide_content(slide, ['●  ' + item for item in leak_items],
    right_col + Inches(0.5), col_top + Inches(1.5), col_w - Inches(1), Inches(3.2),
    font_size=13, color=DARK_TEXT)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 13: Chapter 4 ==========
slide = add_blank_slide()
set_bg(slide, PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5), fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.5), Inches(5), Inches(3), Inches(3), fill_color=SECONDARY)
add_shape(slide, MSO_SHAPE.OVAL, Inches(9.5), Inches(4), Inches(2), Inches(2), fill_color=ACCENT)

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(3), Inches(1), '04', font_size=100, bold=True, color=ACCENT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2), Inches(1), Inches(0.08), fill_color=ACCENT)
add_textbox(slide, Inches(1.5), Inches(4.5), Inches(8), Inches(0.8), '责任体系', font_size=40, bold=True, color=WHITE)
add_textbox(slide, Inches(1.5), Inches(5.3), Inches(8), Inches(0.5), '保密工作组织架构 · 全员参与，各负其责', font_size=18, color=RGBColor(0xB0, 0xC4, 0xDE))

# ========== Slide 14: Responsibility System ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '保密安全工作责任体系', 'Responsibility System')
page = next_page()

# Hierarchy diagram
levels = [
    ('保密安全工作委员会', '最高决策层', PRIMARY, Inches(4.167), Inches(1.7), Inches(5.0), Inches(0.9)),
    ('保密安全办公室', '执行管理部门', SECONDARY, Inches(4.167), Inches(3.0), Inches(5.0), Inches(0.8)),
    ('涉密部门', '具体落实部门', ACCENT, Inches(2.5), Inches(4.3), Inches(3.5), Inches(0.75)),
    ('员工', '直接责任人', RGBColor(0x6C, 0x75, 0x7D), Inches(7.3), Inches(4.3), Inches(3.5), Inches(0.75)),
]

for title, subtitle, color, left, top, w, h in levels:
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill_color=color)
    card.adjustments[0] = 0.1
    add_textbox(slide, left, top + Inches(0.1), w, Inches(0.4), title, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.45), w, Inches(0.3), subtitle, font_size=11, color=RGBColor(0xE0, 0xE8, 0xF0), alignment=PP_ALIGN.CENTER)

# Connecting lines
def add_line(slide, x1, y1, x2, y2, color):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2)

add_line(slide, Inches(6.667), Inches(2.6), Inches(6.667), Inches(3.0), SECONDARY)
add_line(slide, Inches(6.667), Inches(3.8), Inches(6.667), Inches(4.05), LIGHT_GRAY)
add_line(slide, Inches(4.25), Inches(4.05), Inches(6.667), Inches(4.05), LIGHT_GRAY)
add_line(slide, Inches(6.667), Inches(4.05), Inches(9.05), Inches(4.05), LIGHT_GRAY)
add_line(slide, Inches(4.25), Inches(4.05), Inches(4.25), Inches(4.3), ACCENT)
add_line(slide, Inches(9.05), Inches(4.05), Inches(9.05), Inches(4.3), RGBColor(0x6C, 0x75, 0x7D))

add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.4), '员工保密守则', font_size=18, bold=True, color=PRIMARY)

rules = [
    '不该说的秘密不说',
    '不该问的秘密不问',
    '不该看的秘密不看',
    '不该记录的秘密不记录',
    '不在私人交往中谈论秘密',
    '不擅自携带秘密载体外出',
    '不擅自复制、保存和销毁秘密载体',
    '不在无保障的地方存放或使用秘密载体',
    '不在无保障的计算机系统上阅读、存储、传递秘密事项',
]

# 3 columns of rules
for i, rule in enumerate(rules):
    col = i // 3
    row = i % 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(6.0 + row * 0.4)
    
    dot = add_shape(slide, MSO_SHAPE.OVAL, left, top + Inches(0.08), Inches(0.15), Inches(0.15), fill_color=ACCENT)
    add_textbox(slide, left + Inches(0.3), top, Inches(3.7), Inches(0.35), rule, font_size=12, color=DARK_TEXT)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 15: Chapter 5 ==========
slide = add_blank_slide()
set_bg(slide, PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5), fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.5), Inches(5), Inches(3), Inches(3), fill_color=SECONDARY)
add_shape(slide, MSO_SHAPE.OVAL, Inches(9.5), Inches(4), Inches(2), Inches(2), fill_color=ACCENT)

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(3), Inches(1), '05', font_size=100, bold=True, color=ACCENT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2), Inches(1), Inches(0.08), fill_color=ACCENT)
add_textbox(slide, Inches(1.5), Inches(4.5), Inches(8), Inches(0.8), '公司制度', font_size=40, bold=True, color=WHITE)
add_textbox(slide, Inches(1.5), Inches(5.3), Inches(8), Inches(0.5), '保密管理与处罚 · 有章可循，违章必究', font_size=18, color=RGBColor(0xB0, 0xC4, 0xDE))

# ========== Slide 16: Classification Table ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '公司保密范围与密级划分', 'Confidentiality Classification')
page = next_page()

# Table header
table_left = Inches(0.8)
table_top = Inches(1.7)
table_w = Inches(11.7)
col1_w = Inches(1.2)
col2_w = Inches(8.5)
col3_w = Inches(2.0)

header_h = Inches(0.6)
add_shape(slide, MSO_SHAPE.RECTANGLE, table_left, table_top, col1_w, header_h, fill_color=PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, table_left + col1_w, table_top, col2_w, header_h, fill_color=PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, table_left + col1_w + col2_w, table_top, col3_w, header_h, fill_color=PRIMARY)

add_textbox(slide, table_left, table_top + Inches(0.12), col1_w, Inches(0.4), '序号', font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, table_left + col1_w, table_top + Inches(0.12), col2_w, Inches(0.4), '保密范围', font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, table_left + col1_w + col2_w, table_top + Inches(0.12), col3_w, Inches(0.4), '商业密级', font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

rows = [
    ('1', '电子浆料配方（熔断器熔丝浆料配方、灭弧浆料配方；静电抑制器功能、包封浆料配方；化学油墨、钨浆、LTCC电子浆料等配方）', '核心商密', RGBColor(0xFF, 0xEB, 0xEE)),
    ('2', '公司总体经营战略、重要商务谈判内容及载体、正式合同和协议文件', '核心商密', RGBColor(0xFF, 0xEB, 0xEE)),
    ('3', '军工产品设计图纸、生产工艺、工艺参数表、原材料检验标准、作业和检验指导书', '核心商密', RGBColor(0xFF, 0xEB, 0xEE)),
    ('4', '自主研发产品制造工艺、工艺参数表、原材料检验标准、作业和检验指导书', '普通商密', RGBColor(0xE3, 0xF2, 0xFD)),
    ('5', '重要客户及供应商信息', '普通商密', RGBColor(0xE3, 0xF2, 0xFD)),
]

row_h = Inches(0.75)
for i, (num, scope, level, bg_color) in enumerate(rows):
    row_top = table_top + header_h + Inches(i * 0.75)
    add_shape(slide, MSO_SHAPE.RECTANGLE, table_left, row_top, col1_w, row_h, fill_color=bg_color)
    add_shape(slide, MSO_SHAPE.RECTANGLE, table_left + col1_w, row_top, col2_w, row_h, fill_color=bg_color)
    add_shape(slide, MSO_SHAPE.RECTANGLE, table_left + col1_w + col2_w, row_top, col3_w, row_h, fill_color=bg_color)
    
    add_textbox(slide, table_left, row_top + Inches(0.2), col1_w, Inches(0.4), num, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, table_left + col1_w + Inches(0.2), row_top + Inches(0.1), col2_w - Inches(0.4), Inches(0.6), scope, font_size=11, color=DARK_TEXT)
    
    level_color = RGBColor(0xC6, 0x28, 0x28) if level == '核心商密' else RGBColor(0x02, 0x77, 0xBD)
    add_textbox(slide, table_left + col1_w + col2_w, row_top + Inches(0.2), col3_w, Inches(0.4), level, font_size=13, bold=True, color=level_color, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.3),
    '注：已定密为国家秘密的内容按国家秘密进行保护',
    font_size=11, color=GRAY_TEXT, alignment=PP_ALIGN.RIGHT)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 17: Punishment Rules ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '公司对泄密的处罚规定', 'Punishment for Breach of Confidentiality')
page = next_page()

add_textbox(slide, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.5),
    '依据《商业秘密管理办法》第四十二条，有下列行为之一的，根据情节轻重按《员工手册》给予纪律处分，',
    font_size=13, color=DARK_TEXT)
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.5),
    '并可要求承担经济损失赔偿责任；涉嫌违法犯罪的，移送司法机关处理：',
    font_size=13, color=DARK_TEXT)

punishments = [
    ('违反管理规定', '造成商业秘密文件、资料被盗、丢失的'),
    ('泄露造成损害', '泄露商业秘密已造成损害后果的'),
    ('强制他人违规', '利用职权强制他人违反保密规定造成泄密的'),
    ('提供便利条件', '为他人泄露或侵犯公司商业秘密提供便利条件的'),
    ('非法出卖牟利', '以谋取私利为目的，非法出卖和泄露商业秘密的'),
    ('其他违规情形', '其他因玩忽职守、违反规定造成公司商业秘密泄露的'),
]

for i, (title, desc) in enumerate(punishments):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(2.8 + row * 2.1)
    w = Inches(3.7)
    h = Inches(1.8)
    
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill_color=WHITE)
    card.adjustments[0] = 0.05
    
    num_shape = add_shape(slide, MSO_SHAPE.OVAL, left + Inches(0.25), top + Inches(0.2), Inches(0.55), Inches(0.55), fill_color=RGBColor(0xDC, 0x35, 0x45))
    add_textbox(slide, left + Inches(0.25), top + Inches(0.28), Inches(0.55), Inches(0.4), str(i+1), font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, left + Inches(0.95), top + Inches(0.25), w - Inches(1.2), Inches(0.45), title, font_size=15, bold=True, color=PRIMARY)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.9), w - Inches(0.5), Inches(0.8), desc, font_size=11, color=GRAY_TEXT)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 18: Chapter 6 + Protection Measures ==========
slide = add_blank_slide()
set_bg(slide, PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5), fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.5), Inches(5), Inches(3), Inches(3), fill_color=SECONDARY)
add_shape(slide, MSO_SHAPE.OVAL, Inches(9.5), Inches(4), Inches(2), Inches(2), fill_color=ACCENT)

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(3), Inches(1), '06', font_size=100, bold=True, color=ACCENT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2), Inches(1), Inches(0.08), fill_color=ACCENT)
add_textbox(slide, Inches(1.5), Inches(4.5), Inches(8), Inches(0.8), '防护措施', font_size=40, bold=True, color=WHITE)
add_textbox(slide, Inches(1.5), Inches(5.3), Inches(8), Inches(0.5), '企业自我保护 · 全方位筑牢保密防线', font_size=18, color=RGBColor(0xB0, 0xC4, 0xDE))

# ========== Slide 19: Protection Measures Detail ==========
slide = add_blank_slide()
set_bg(slide, LIGHT_BG)
add_decorative_header(slide, '企业商业秘密自我保护措施', 'Self-Protection Measures')
page = next_page()

measures = [
    ('门禁制度', '熔断器事业部、浆料实验室为商密保护区域，门禁常闭，须刷卡进入。', '1'),
    ('区域管理', '非熔断器部门人员不得借道进出；外来人员不得进入熔断器事业部、浆料实验室，接待在大厅或会议室。', '2'),
    ('文件管理', '熔断器用工艺文件、作业指导书、检验指导书、质量记录均为核心商密，非授权者不得阅读、抄录、借出，不得带出公司。', '3'),
    ('人员管理', '对从事熔断器和浆料的相关人员签订保密协议，防止竞争对手实施挖人等不正当竞争手段。', '4'),
    ('信息输出管理', '严格实行信息输出审批制度，对刻录结果进行确认，防止核心商密通过互联网泄漏。', '5'),
    ('合作管理', '对外合作、技术交流须经审批，严格控制信息披露范围，签订保密协议后方可进行。', '6'),
]

for i, (title, desc, num) in enumerate(measures):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(1.7 + row * 2.6)
    w = Inches(3.7)
    h = Inches(2.3)
    
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill_color=WHITE)
    card.adjustments[0] = 0.05
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.08), fill_color=ACCENT)
    
    num_bg = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.25), top + Inches(0.3), Inches(0.7), Inches(0.7), fill_color=SECONDARY)
    num_bg.adjustments[0] = 0.3
    add_textbox(slide, left + Inches(0.25), top + Inches(0.38), Inches(0.7), Inches(0.5), num, font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, left + Inches(1.1), top + Inches(0.35), w - Inches(1.3), Inches(0.5), title, font_size=17, bold=True, color=PRIMARY)
    add_textbox(slide, left + Inches(0.3), top + Inches(1.1), w - Inches(0.5), Inches(1.1), desc, font_size=11, color=GRAY_TEXT)

add_footer(slide, page, TOTAL_PAGES)

# ========== Slide 20: End ==========
slide = add_blank_slide()
set_bg(slide, PRIMARY)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5), fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.333), Inches(0.08), fill_color=SECONDARY)
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.5), Inches(0.8), Inches(2.5), Inches(2.5), fill_color=SECONDARY)
add_shape(slide, MSO_SHAPE.OVAL, Inches(11), Inches(1.3), Inches(1.8), Inches(1.8), fill_color=ACCENT)

add_textbox(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(1.2), '感谢聆听', font_size=56, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.917), Inches(4.1), Inches(1.5), Inches(0.08), fill_color=ACCENT)
add_textbox(slide, Inches(0), Inches(4.4), Inches(13.333), Inches(0.6), 'Thank You', font_size=24, color=RGBColor(0xB0, 0xC4, 0xDE), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(0.4), '苏州晶讯科技股份有限公司保密办', font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(6.0), Inches(13.333), Inches(0.4), '保守商业秘密，守护企业未来', font_size=14, color=RGBColor(0xB0, 0xC4, 0xDE), alignment=PP_ALIGN.CENTER)

# Save
output_path = 'D:\\administrator\\Desktop\\商业保密培训_优化版.pptx'
prs.save(output_path)
print(f'Saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
